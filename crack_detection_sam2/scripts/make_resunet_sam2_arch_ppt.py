from __future__ import annotations

import random
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "crack_detection_sam2" / "results" / "architecture"
SLICE_ROOT = ROOT / "_data" / "craq_0-94_v1" / "tiles_512" / "images"
REFERENCE = ROOT / "_literature" / "picture" / "sam2-architecture.png"


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color="#172033"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_block(slide, x, y, w, h, text, fill, line="#8d98a8", size=16, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(1.25)
    tf = shp.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb("#172033")
    return shp


def add_arrow(slide, x1, y1, x2, y2, color="#172033", width=2.0):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_x(slide, x, y, w, h, color="#c7372f"):
    l1 = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + w), Inches(y + h)
    )
    l2 = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x + w), Inches(y), Inches(x), Inches(y + h)
    )
    for line in (l1, l2):
        line.line.color.rgb = rgb(color)
        line.line.width = Pt(2.2)


def make_slice_stack(paths: list[Path], out_path: Path) -> None:
    canvas = Image.new("RGBA", (760, 420), (255, 255, 255, 0))
    offsets = [(0, 44), (62, 28), (124, 12), (186, 0)]
    for idx, (path, (x, y)) in enumerate(zip(paths, offsets)):
        img = Image.open(path).convert("RGB").resize((300, 300))
        shadow = Image.new("RGBA", (316, 316), (0, 0, 0, 0))
        d = ImageDraw.Draw(shadow)
        d.rounded_rectangle((10, 10, 306, 306), radius=16, fill=(0, 0, 0, 58))
        shadow = shadow.filter(ImageFilter.GaussianBlur(9))
        canvas.alpha_composite(shadow, (x + 10, y + 10))
        framed = Image.new("RGBA", (316, 316), (255, 255, 255, 255))
        framed.alpha_composite(img.convert("RGBA"), (8, 8))
        fd = ImageDraw.Draw(framed)
        fd.rounded_rectangle((4, 4, 312, 312), radius=16, outline=(255, 255, 255, 255), width=8)
        fd.rounded_rectangle((8, 8, 308, 308), radius=12, outline=(112, 124, 137, 180), width=2)
        canvas.alpha_composite(framed, (x, y))
        if idx == len(offsets) - 1:
            label = "random 512x512 slices"
            d = ImageDraw.Draw(canvas)
            d.rounded_rectangle((x + 18, y + 258, x + 282, y + 296), radius=12, fill=(15, 23, 42, 190))
            d.text((x + 34, y + 266), label, fill=(255, 255, 255, 255))
    canvas.save(out_path)


def set_slide_bg(slide, color="#f8fafc"):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(color)


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.45, 0.18, 12.4, 0.35, title, size=24, bold=True)
    if subtitle:
        add_textbox(slide, 0.72, 0.55, 11.85, 0.28, subtitle, size=10.5, color="#64748b")


def add_slide_architecture(prs: Presentation, stack_png: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(
        slide,
        "ResUNet-SAM2 Prompt-Refine Architecture",
        "Input is random 512x512 image slices; SAM2 video memory modules are removed from this pipeline.",
    )

    slide.shapes.add_picture(str(stack_png), Inches(0.45), Inches(1.25), width=Inches(2.25))
    add_textbox(slide, 0.44, 4.38, 2.35, 0.35, "RGB tile batch\n[B,3,512,512]", size=11.5, bold=True)

    add_block(slide, 3.0, 1.35, 1.55, 0.82, "ResUNet-50\nstage 1", "#c8f7c5", "#72bf75")
    add_block(slide, 4.85, 1.35, 1.60, 0.82, "craquelure\nprobability map", "#e5f2ff", "#7bb7e8", size=14)
    add_textbox(slide, 4.88, 2.15, 1.55, 0.27, "[B,2,512,512]", size=9.5, color="#64748b")

    add_block(slide, 3.0, 3.35, 1.65, 0.82, "SAM2 image\nencoder", "#d7f1ff", "#67a9cf")
    add_textbox(slide, 3.02, 4.18, 1.62, 0.25, "frozen Hiera", size=9.5, color="#64748b")
    add_block(slide, 6.85, 2.55, 1.65, 0.82, "SAM2 prompt\nencoder", "#d9d7ff", "#918ce0", size=14)
    add_block(slide, 8.95, 2.55, 1.65, 0.82, "SAM2 mask\ndecoder", "#ffe2bf", "#d8964f", size=14)
    add_textbox(slide, 6.8, 3.42, 1.75, 0.23, "trainable", size=9.5, color="#64748b")
    add_textbox(slide, 8.9, 3.42, 1.75, 0.23, "trainable", size=9.5, color="#64748b")

    add_block(slide, 11.15, 2.55, 1.45, 0.82, "refined\nmask", "#ffd6dd", "#cc6a7a", size=15)
    add_textbox(slide, 11.15, 3.42, 1.45, 0.23, "[B,1,512,512]", size=9.5, color="#64748b")

    add_arrow(slide, 2.54, 2.07, 3.0, 1.77)
    add_arrow(slide, 4.55, 1.76, 4.85, 1.76)
    add_arrow(slide, 2.54, 2.58, 3.0, 3.76)
    add_arrow(slide, 4.65, 3.76, 8.95, 2.96, color="#64748b", width=1.6)
    add_arrow(slide, 6.45, 1.76, 6.85, 2.96, color="#0f766e")
    add_arrow(slide, 8.50, 2.96, 8.95, 2.96)
    add_arrow(slide, 10.60, 2.96, 11.15, 2.96)

    add_textbox(slide, 5.95, 1.02, 1.55, 0.26, "dense mask prompt\n(resized to 128/256)", size=9.5, color="#0f766e")
    add_textbox(slide, 5.15, 3.95, 2.55, 0.27, "image embeddings + optional high-res features", size=9.5, color="#64748b")

    removed = add_block(slide, 7.0, 4.52, 5.35, 0.78, "Removed / not used: memory attention  |  memory encoder  |  memory bank", "#eeeeee", "#b7bdc7", size=13)
    removed.text_frame.paragraphs[0].runs[0].font.color.rgb = rgb("#475569")
    add_x(slide, 7.12, 4.58, 0.55, 0.58)
    add_x(slide, 9.03, 4.58, 0.55, 0.58)
    add_x(slide, 10.82, 4.58, 0.55, 0.58)

    add_textbox(slide, 0.58, 6.83, 12.1, 0.28, "Implementation basis: PromptedSAM2Seg keeps only image_encoder, sam_prompt_encoder, and sam_mask_decoder from the built SAM2 model.", size=10.5, color="#475569")


def add_slide_reference(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Modification From SAM2 Video Diagram", "The video-memory branch in the reference architecture is excluded for ResUNet-SAM2.")
    slide.shapes.add_picture(str(REFERENCE), Inches(0.55), Inches(1.0), width=Inches(12.25))
    add_block(slide, 4.75, 3.95, 3.85, 0.55, "kept: image encoder + prompt encoder + mask decoder", "#e7f5ff", "#74a9cf", size=13)
    add_block(slide, 8.92, 4.75, 3.45, 0.55, "removed: memory encoder / bank", "#eeeeee", "#b7bdc7", size=13)
    add_block(slide, 4.24, 1.18, 1.6, 0.5, "removed:\nmemory attention", "#eeeeee", "#b7bdc7", size=11)
    add_x(slide, 4.27, 1.2, 1.54, 0.46)
    add_x(slide, 9.0, 4.8, 3.25, 0.45)


def add_slide_verification(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Module Verification", "Checked against local implementation files.")
    rows = [
        ("Used", "crack_detection_sam2/model_prompted_sam2.py", "self.image_encoder = sam2.image_encoder"),
        ("Used", "crack_detection_sam2/model_prompted_sam2.py", "self.sam_prompt_encoder = sam2.sam_prompt_encoder"),
        ("Used", "crack_detection_sam2/model_prompted_sam2.py", "self.sam_mask_decoder = sam2.sam_mask_decoder"),
        ("Not used", "crack_detection_sam2/model_prompted_sam2.py", "no memory_attention / memory_encoder / memory_bank assignment or forward call"),
        ("Not relevant", "crack_detection_sam2/train_craq_promptrefine.py", "only memory hit is DataLoader pin_memory"),
    ]
    x0, y0 = 0.9, 1.22
    widths = [1.2, 4.25, 6.1]
    headers = ["Status", "File", "Evidence"]
    for i, h in enumerate(headers):
        add_block(slide, x0 + sum(widths[:i]), y0, widths[i], 0.42, h, "#dbeafe", "#93c5fd", size=12)
    for r, row in enumerate(rows):
        y = y0 + 0.5 + r * 0.66
        fill = "#ffffff" if r % 2 == 0 else "#f1f5f9"
        for i, val in enumerate(row):
            shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0 + sum(widths[:i])), Inches(y), Inches(widths[i]), Inches(0.56))
            shp.fill.solid()
            shp.fill.fore_color.rgb = rgb(fill)
            shp.line.color.rgb = rgb("#cbd5e1")
            tf = shp.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = val
            run.font.size = Pt(10.5)
            run.font.color.rgb = rgb("#172033" if row[0] == "Used" else "#475569")
    add_textbox(slide, 0.9, 5.35, 11.55, 0.6, "Conclusion: the architecture is an image-segmentation prompt-refine pipeline, not the full SAM2 video architecture. The memory path should be removed from the presentation diagram.", size=14, bold=True, color="#0f172a")


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    random.seed(512)
    candidates = sorted(SLICE_ROOT.glob("*.png"))
    if len(candidates) < 4:
        raise RuntimeError(f"Need at least 4 slice images under {SLICE_ROOT}")
    chosen = random.sample(candidates, 4)
    stack_png = OUT_DIR / "random_512_slices_stack.png"
    make_slice_stack(chosen, stack_png)

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    add_slide_architecture(prs, stack_png)
    add_slide_reference(prs)
    add_slide_verification(prs)

    out = OUT_DIR / "resunet_sam2_architecture.pptx"
    prs.save(out)
    print(out)
    print("Selected slices:")
    for p in chosen:
        print(p)


if __name__ == "__main__":
    main()
